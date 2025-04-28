# knowledge/loader.py
import re
import csv
from pathlib import Path
from typing import List, Optional
from PIL import Image
import fitz  # PyMuPDF for PDF
import pytesseract  # OCR for scanned PDFs
from docx import Document as DocxDocument  # python-docx for .docx
from pptx import Presentation  # python-pptx for .pptx slides
from langchain.docstore.document import Document
from langchain_community.vectorstores import FAISS
from backend.app.experts.embedder import get_embedding
from utils.config import DOCS_PATH, VECTORSTORE_PATH
from bs4 import BeautifulSoup
import pandas as pd
import matplotlib.pyplot as plt
import io, base64
import json





class OllamaEmbeddings:
    """Embedding interface for Ollama models"""
    embedding_dimension: int = 4096

    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        return [get_embedding(t) for t in texts]

    def embed_query(self, text: str) -> List[float]:
        return get_embedding(text)
# --- Algorithm Extraction ---
def extract_algorithms(text: str) -> List[str]:
    """Extract algorithm names supporting word, hyphens, underscores."""
    patterns = [
        r"Algorithm[:\s]+([\w\-]+)",
        r"([\w\-]+)\s+algorithm",
        r"(?:procedure|method)[:\s]*([\w\-]+)"
    ]
    algos = set()
    for p in patterns:
        for m in re.findall(p, text, flags=re.IGNORECASE):
            algos.add(m.strip())
    return list(algos)

# --- File Loaders ---
def load_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")

def load_csv_file(path: Path) -> str:
    rows = []
    with path.open(newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(", ".join(row))
    return "\n".join(rows)

def load_pdf_documents(path: Path) -> List[Document]:
    """Load a PDF and return as a single Document"""
    text = load_pdf_file(path)
    return [Document(page_content=text, metadata={"source": str(path)})]

def load_pdf_file(path: Path) -> str:
    text_chunks: List[str] = []
    pdf = fitz.open(str(path))
    for page in pdf:
        txt = page.get_text().strip()
        if txt:
            text_chunks.append(txt)
        else:
            pix = page.get_pixmap(dpi=300)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            ocr = pytesseract.image_to_string(img)
            text_chunks.append(ocr)
    return "\n".join(text_chunks)

def load_docx_file(path: Path) -> str:
    doc = DocxDocument(str(path))
    texts: List[str] = [p.text for p in doc.paragraphs if p.text]
    for table in doc.tables:
        for row in table.rows:
            texts.append(", ".join(cell.text for cell in row.cells))
    return "\n".join(texts)

def load_pptx_file(path: Path) -> str:
    prs = Presentation(str(path))
    texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)

def load_file(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in [".txt", ".md"]:
        return load_text_file(path)
    if ext == ".csv":
        return load_csv_file(path)
    if ext == ".pdf":
        return load_pdf_file(path)
    if ext == ".docx":
        return load_docx_file(path)
    if ext in [".pptx", ".ppt"]:
        return load_pptx_file(path)
    raise ValueError(f"Unsupported file type: {ext}")



def extract_media_from_html(html: str, source_path: Path):
    """
    فرض می‌کنیم اسناد HTML یا Markdown حاوی <img>, <table> و <canvas data-chart> هستند.
    """
    soup = BeautifulSoup(html, "html.parser")
    media = []

    # تصاویر
    for img in soup.find_all("img"):
        src = img.get("src")
        # مسیر نسبی به مسیری روی سرور
        url = f"/media/{source_path.stem}/{src}"
        media.append({"type": "image", "url": url, "caption": img.get("alt", "")})

    # جداول
    for idx, table in enumerate(soup.find_all("table")):
        df = pd.read_html(str(table))[0]
        # ذخیره JSON یا HTML
        media_id = f"{source_path.stem}-table-{idx}"
        out_json = source_path.parent / "media" / f"{media_id}.json"
        out_json.parent.mkdir(exist_ok=True, parents=True)
        df.to_json(out_json, orient="records")
        media.append({"type": "table", "url": f"/media/{out_json.name}", "caption": ""})

    # نمودارها (مثلا در مارک‌داون <canvas data-chart="...">)
    for idx, canvas in enumerate(soup.find_all("canvas", {"data-chart": True})):
        chart_spec = json.loads(canvas["data-chart"])
        # می‌توانید با matplotlib دوباره نمودار را رسم و ذخیره کنید:
        fig, ax = plt.subplots()
        ax.plot(chart_spec["x"], chart_spec["y"])
        buf = io.BytesIO()
        fig.savefig(buf, format="png")
        buf.seek(0)
        img_bytes = buf.read()
        b64 = base64.b64encode(img_bytes).decode()
        media_id = f"{source_path.stem}-chart-{idx}.png"
        out_img = source_path.parent / "media" / media_id
        out_img.parent.mkdir(exist_ok=True, parents=True)
        out_img.write_bytes(img_bytes)
        media.append({"type": "chart", "url": f"/media/{media_id}", "caption": canvas.get("data-caption", "")})
        plt.close(fig)

    return media


# --- Vectorstore Builder ---
def build_vectorstore(source_dir: Optional[str] = None):
    src = Path(source_dir) if source_dir else DOCS_PATH
    src.mkdir(parents=True, exist_ok=True)
    VECTORSTORE_PATH.mkdir(parents=True, exist_ok=True)

    # gather all supported files
    patterns = ["*.md", "*.txt", "*.csv", "*.pdf", "*.docx", "*.pptx", "*.ppt"]
    files = [p for pattern in patterns for p in src.rglob(pattern)]

    documents, metadatas = [], []
    all_algos = []

    for fp in files:
        content = load_file(fp)
        documents.append(Document(page_content=content, metadata={"source": str(fp)}))
        metadatas.append({"source": str(fp), "algorithms": extract_algorithms(content)})
        all_algos.extend(extract_algorithms(content))

    if not documents:
        print("No documents to index.")
        return

    # index text docs
    embeddings_fn = get_embedding
    db = FAISS.from_documents(documents, embeddings_fn, metadatas=metadatas)
    db.save_local(str(VECTORSTORE_PATH))
    print(f"Built text vectorstore with {len(documents)} docs.")

    # index unique algorithms
    unique_algos = set(all_algos)
    if unique_algos:
        algo_docs = [Document(page_content=a, metadata={}) for a in unique_algos]
        algo_dir = VECTORSTORE_PATH / "algos"
        algo_dir.mkdir(parents=True, exist_ok=True)
        algo_db = FAISS.from_documents(algo_docs, embeddings_fn)
        algo_db.save_local(str(algo_dir))
        print(f"Built algorithm index with {len(unique_algos)} algos.")


def search_knowledge(query: str, k: int = 5) -> Optional[List[str]]:
    """Retrieve top-k doc contents from FAISS."""
    if not VECTORSTORE_PATH.exists():
        return []
    db = FAISS.load_local(str(VECTORSTORE_PATH), OllamaEmbeddings().embed_query)
    results = db.similarity_search(query, k=k)
    return [doc.page_content for doc in results]

def search_algorithms(query: str, k: int = 5) -> List[str]:
    """
    Retrieve top-k algorithm names semantically similar to the query.
    """
    algo_dir = VECTORSTORE_DIR / "algos"
    if not algo_dir.exists():
        return []
    db = FAISS.load_local(str(algo_dir), embedding_fn)
    results = db.similarity_search(query, k=k)
    return [doc.page_content for doc in results]


def retrieve_with_media(query: str, k: int = 5):
    # جستجوی متن
    texts = search_knowledge(query, k)
    # بارگذاری ایندکس رسانه
    media_index = json.loads((VECTORSTORE_PATH / "media_index.json").read_text(encoding="utf-8"))
    # فیلتر رسانه‌های مرتبط با نتایج
    related_media = []
    sources = {doc.metadata["source"] for doc in texts}
    for m in media_index:
        if m["source"] in sources:
            related_media.append(m)
    return {"texts": [d.page_content for d in texts], "media": related_media}




if __name__ == "__main__":
    build_vectorstore()